import os
import re
from typing import List, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from .base_converter import BaseConverter

class MdToPptConverter(BaseConverter):
    """
    Converts Markdown files to PowerPoint (.pptx) presentations.
    This class encapsulates the logic from the original tools/md_to_ppt.py script.
    """

    def __init__(self, output_dir: str, **kwargs):
        super().__init__(output_dir, **kwargs)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Initialize attributes from kwargs, similar to MdToOfficeConverter
        self.template_path = kwargs.get('template_path')
        self.project_name = kwargs.get('project_name', '')
        self.author = kwargs.get('author', '')
        # Add any other relevant attributes from kwargs
        
        # Logic to find default template if not provided
        if not self.template_path:
            # This logic can be adapted from the original script or MdToOfficeConverter
            default_template = Path(__file__).parent / "templates" / "template.pptx"
            if default_template.exists():
                self.template_path = str(default_template)
                self.logger.info(f"Using default PPTX template: {self.template_path}")
            else:
                self.logger.warning("Default PPTX template not found. A blank presentation will be created.")

    def convert(self, input_path: str) -> List[str]:
        """
        Main conversion entry point. Handles both single files and directories.
        """
        if not self._is_valid_input(input_path, ['.md']):
            raise ValueError(f"Invalid input file for PPTX conversion: {input_path}")

        output_files = []
        if os.path.isfile(input_path):
            output_file = self._convert_single_file(input_path)
            if output_file:
                output_files.append(output_file)
        else:
            md_files = self._get_files_by_extension(input_path, ['.md'])
            for md_file in md_files:
                output_file = self._convert_single_file(md_file)
                if output_file:
                    output_files.append(output_file)
        
        return output_files

    def _convert_single_file(self, md_file: str) -> Optional[str]:
        """
        Processes a single Markdown file and creates a presentation.
        This is where the main logic from `process_md_file` will go.
        """
        output_file = self.output_dir / f"{Path(md_file).stem}.pptx"
        self.logger.info(f"Starting conversion: {md_file} -> {output_file}")
        
        try:
            # 1. Parse the markdown content
            sections = self._parse_markdown(md_file)
            if not sections:
                self.logger.warning(f"No content found in {md_file}. Skipping.")
                return None
            
            # 2. Create presentation
            if self.template_path and Path(self.template_path).exists():
                try:
                    prs = Presentation(self.template_path)
                    # Clear existing slides from template
                    for i in range(len(prs.slides) - 1, -1, -1):
                        rId = prs.slides._sldIdLst[i].rId
                        prs.part.drop_rel(rId)
                        del prs.slides._sldIdLst[i]
                except Exception as e:
                    self.logger.error(f"Failed to load template {self.template_path}: {str(e)}")
                    self.logger.info("Creating blank presentation as fallback")
                    prs = Presentation()
            else:
                prs = Presentation()

            # 3. Create slides based on sections
            # Create title slide for the first section (level 0)
            if sections:
                title_section = sections.pop(0)
                self._create_title_slide(prs, title_section['title'])

            # Create content slides for the rest
            for section in sections:
                try:
                    # Split content into multiple slides if necessary
                    content_slides_data = self._split_content_into_slides(section['content'])
                    
                    is_first_slide_of_section = True
                    for slide_content_lines in content_slides_data:
                        # For the first slide of a new section, use the section title.
                        # For subsequent slides of the same section, use a continuation title.
                        slide_title = section['title'] if is_first_slide_of_section else f"{section['title']} (续)"
                        self._create_content_slide(prs, slide_title, slide_content_lines)
                        is_first_slide_of_section = False
                except Exception as e:
                    self.logger.error(f"Failed to create slide for section '{section['title']}': {str(e)}")
                    # Continue with next section instead of failing entire conversion
                    continue

            # 4. Save the presentation
            try:
                self._save_presentation_with_retry(prs, str(output_file))
                self.logger.info(f"Successfully converted {md_file} to {output_file}")
                return str(output_file)
            except Exception as e:
                self.logger.error(f"Failed to save presentation: {str(e)}")
                if isinstance(e, PermissionError):
                    raise ValueError(f"无法保存文件 {output_file}，文件可能正在被其他程序使用。")
                raise ValueError(f"保存演示文稿时出错：{str(e)}")
            
        except Exception as e:
            error_msg = f"转换 {md_file} 时出错：{str(e)}"
            self.logger.error(error_msg)
            self.logger.debug(traceback.format_exc())
            raise ValueError(error_msg)

    def _parse_markdown(self, md_file: str) -> List[dict]:
        """
        Parses a Markdown file and returns a structured list of sections.
        Migrated from tools/md_to_ppt.py
        """
        with open(md_file, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Default title from filename
        title = Path(md_file).stem
        
        # Try to get title from the first H1 heading
        first_h1_match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
        if first_h1_match:
            title = first_h1_match.group(1).strip()
        
        sections = []
        # The first "section" is the cover page title
        sections.append({'level': 0, 'title': title, 'content': []})
        
        # Split content by H2 and lower headings
        # Using a more robust split pattern to handle various heading styles
        parts = re.split(r'^(##+)\s+(.+?)\s*$', content, flags=re.MULTILINE)
        
        # The first part is the content before the first real heading
        # This is often empty or intro text, we can associate it with the title slide if needed
        
        # Process the remaining parts
        i = 1
        while i < len(parts):
            level = len(parts[i]) # Number of '#' gives the heading level
            heading_text = parts[i+1].strip()
            section_content = parts[i+2].strip().split('\n')
            
            sections.append({'level': level, 'title': heading_text, 'content': section_content})
            i += 3
            
        return sections

    def _split_content_into_slides(self, content_lines: list, max_points: int = 8, max_len: int = 500) -> list:
        """
        Splits a list of content lines into multiple lists, each for one slide.
        Migrated and simplified from tools/md_to_ppt.py.
        """
        if not content_lines:
            return [[]]

        slides = []
        current_slide = []
        point_count = 0
        char_count = 0

        for line in content_lines:
            stripped_line = line.strip()
            if not stripped_line:
                continue

            line_len = len(stripped_line)
            is_bullet = stripped_line.startswith('-') or stripped_line.startswith('*') or re.match(r'^\d+\.', stripped_line)

            # Check if adding this line would exceed the slide's capacity
            if current_slide and ( (is_bullet and point_count >= max_points) or (char_count + line_len > max_len) ):
                slides.append(current_slide)
                current_slide = []
                point_count = 0
                char_count = 0

            current_slide.append(line)
            char_count += line_len
            if is_bullet:
                point_count += 1
        
        if current_slide:
            slides.append(current_slide)
            
        return slides

    def _find_best_layout(self, prs: Presentation, layout_type: str):
        """Finds the best slide layout based on type ('title' or 'content')."""
        try:
            # Simple strategy: 0 for title, 1 for section header, 5 for content
            if layout_type == 'title':
                return prs.slide_layouts[0]
            else: # 'content'
                # In many standard templates, layout 1 is 'Title and Content'
                # and layout 5 is 'Blank'. We prefer 'Title and Content'.
                try:
                    return prs.slide_layouts[1] # 'Title and Content' is often at index 1
                except IndexError:
                    return prs.slide_layouts[0] # Fallback to the first available layout
        except Exception as e:
            self.logger.error(f"Failed to find layout for type '{layout_type}': {str(e)}")
            raise ValueError(f"无法找到合适的幻灯片布局：{str(e)}")

    def _create_title_slide(self, prs: Presentation, title: str):
        """Creates a title slide."""
        try:
            layout = self._find_best_layout(prs, 'title')
            slide = prs.slides.add_slide(layout)
            
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
            
            # Optionally add author/project name from self.author etc.
            # to a subtitle placeholder if it exists
            try:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = self.author or self.project_name
            except (KeyError, IndexError):
                pass # No subtitle placeholder
        except Exception as e:
            self.logger.error(f"Failed to create title slide: {str(e)}")
            raise ValueError(f"创建标题页时出错：{str(e)}")

    def _create_content_slide(self, prs: Presentation, title: str, content: list):
        """Creates a content slide."""
        try:
            layout = self._find_best_layout(prs, 'content')
            slide = prs.slides.add_slide(layout)

            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title

            # Find the body/content placeholder
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    body_shape = shape
                    break
            
            if body_shape:
                tf = body_shape.text_frame
                tf.clear() # Clear default text
                tf.word_wrap = True
                
                for line in content:
                    p = tf.add_paragraph()
                    p.text = line.strip()
                    # Basic indentation for lists
                    if p.text.startswith('- ') or p.text.startswith('* '):
                        p.text = p.text[2:]
                        p.level = 1
                    elif re.match(r'^\d+\.\s', p.text):
                        p.text = re.sub(r'^\d+\.\s', '', p.text)
                        p.level = 1
                    else:
                        p.level = 0
            else:
                self.logger.warning("No body placeholder found on content slide layout.")
        except Exception as e:
            self.logger.error(f"Failed to create content slide: {str(e)}")
            raise ValueError(f"创建内容页时出错：{str(e)}")

    def _save_presentation_with_retry(self, prs: Presentation, output_file: str, max_retries: int = 3, delay: int = 1) -> bool:
        """
        Saves the presentation with a retry mechanism.
        Migrated from tools/md_to_ppt.py
        """
        import time
        
        Path(output_file).parent.mkdir(parents=True, exist_ok=True)
        
        for attempt in range(max_retries):
            try:
                # If file exists, try to remove it first
                if os.path.exists(output_file):
                    try:
                        os.remove(output_file)
                    except PermissionError:
                        if attempt < max_retries - 1:
                            self.logger.warning(f"File {output_file} is in use. Retrying in {delay}s...")
                            time.sleep(delay)
                            continue
                        else:
                            raise ValueError(f"文件 {output_file} 正在被其他程序使用，无法保存。")
                
                prs.save(output_file)
                self.logger.info(f"Successfully saved {output_file}")
                return True
                
            except Exception as e:
                if attempt < max_retries - 1:
                    self.logger.warning(f"Save attempt {attempt + 1} failed: {str(e)}. Retrying...")
                    time.sleep(delay)
                else:
                    raise ValueError(f"保存文件失败：{str(e)}")
        
        return False
